"""
Document Processor Service - ENHANCED VERSION
- Preserves formatting as markdown
- Extracts images from PDFs
- Generates CLIP embeddings for images
"""

import io
import logging
from typing import Dict, Any, Optional, List
from datetime import datetime
import json
import uuid
import os
from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
from PIL import Image

# NEW: For PDF image extraction
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logging.warning("PyMuPDF not installed - PDF image extraction disabled")

from app.config import settings
from app.services.text_processor import TextProcessor
from app.services.embeddings import EmbeddingService
from app.services.vector_store import VectorStore
from app.services.pdf_image_extractor import PDFImageExtractor
from app.models.responses import DocumentProcessingResult, EmbeddingResult

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process documents with formatting preservation and image extraction"""
    
    def __init__(self):
        self.text_processor = TextProcessor()
        self.embedding_service = EmbeddingService()
        self.vector_store = VectorStore()
        self.pdf_image_extractor = PDFImageExtractor()
        
        # NEW: Check if image processor available
        try:
            from app.services.image_processor import ImageProcessor
            self.image_processor = ImageProcessor()
            logger.info("Image processor initialized")
        except Exception as e:
            self.image_processor = None
            logger.warning(f"Image processor not available: {e}")
    
    async def _create_document_record(
        self,
        document_id: str,
        kb_id: str,
        tenant_id: str,
        filename: str,
        file_size: int,
        content_type: str,
        metadata: Dict[str, Any]
    ):
        """Create document record in database"""
        import mysql.connector
        from app.config import settings
        
        conn = mysql.connector.connect(
            host=settings.DB_HOST,
            port=settings.DB_PORT,
            user=settings.DB_USER,
            password=settings.DB_PASSWORD,
            database=settings.DB_NAME
        )
        cursor = conn.cursor()
        
        try:
            cursor.execute("""
                INSERT INTO yovo_tbl_aiva_documents 
                (id, kb_id, tenant_id, filename, original_filename, file_type, 
                 file_size_bytes, storage_url, status, metadata, created_at, updated_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, NOW(), NOW())
            """, (
                document_id,
                kb_id,
                tenant_id,
                filename,
                filename,
                content_type,
                file_size,
                f"/storage/documents/{document_id}",
                "processing",
                json.dumps(metadata)
            ))
            
            conn.commit()
            logger.info(f"Created document record: {document_id}")
            
        except Exception as e:
            conn.rollback()
            logger.error(f"Error creating document record: {e}")
            raise
        finally:
            cursor.close()
            conn.close()

    async def process_document(
        self,
        document_id: str,
        kb_id: str,
        tenant_id: str,
        file_content: bytes,
        filename: str,
        content_type: str,
        metadata: Dict[str, Any]
    ) -> Any:
        """Process document with markdown preservation and image extraction"""
        import time
        start_time = time.time()
        
        logger.info(f"Processing document: {filename} ({content_type})")
        
        # Create document record
        #await self._create_document_record(
        #    document_id=document_id,
        #    kb_id=kb_id,
        #    tenant_id=tenant_id,
        #    filename=filename,
        #    file_size=len(file_content),
        #    content_type=content_type,
        #    metadata=metadata
        #)
        
        # NEW: Extract text AND images
        extraction_result = await self._extract_content(
            file_content, 
            filename, 
            content_type,
            document_id=document_id,  # ✅ Added
            kb_id=kb_id,              # ✅ Added
            tenant_id=tenant_id       # ✅ Added
        )
        
        # Process text with markdown preservation
        processed = await self.text_processor.process_text(
            text=extraction_result["text"],
            document_id=document_id,
            kb_id=kb_id,
            metadata={
                **metadata,
                "filename": filename,
                "content_type": content_type,
                "pages": extraction_result.get("pages", 0),
                "extracted_images": extraction_result.get("images", 0)  # ✅ Added
            },
            preserve_formatting=True  # NEW: Preserve markdown
        )
        
        # Generate embeddings for text chunks
        embeddings_result = await self.embedding_service.generate_embeddings_for_chunks(
            chunks=processed["chunks"]
        )
        
        # Store in vector store
        await self.vector_store.store_document(
            document_id=document_id,
            kb_id=kb_id,
            tenant_id=tenant_id,
            chunks=processed["chunks"],
            embeddings=embeddings_result["embeddings"]
        )
        
        processing_time = int((time.time() - start_time) * 1000)
        
        # Build response
        processing_results = DocumentProcessingResult(
            total_pages=extraction_result.get("pages", 0),
            total_chunks=len(processed["chunks"]),
            extracted_images=extraction_result.get("images", 0),  # UPDATED: Real count
            detected_tables=extraction_result.get("tables", 0),
            detected_faqs=processed.get("faqs", 0),
            chunks_by_type=processed.get("chunks_by_type", {}),
            language_detected=processed.get("languages", []),
            has_roman_urdu=processed.get("has_roman_urdu", False),
            processing_time_ms=processing_time
        )
        
        embedding_result = EmbeddingResult(
            total_embeddings_generated=embeddings_result["total_embeddings"],
            embedding_model=embeddings_result["model"],
            total_tokens_embedded=embeddings_result["total_tokens"],
            vector_dimension=embeddings_result["dimension"]
        )
        
        class ProcessingResponse:
            def __init__(self, proc, emb):
                self.processing_results = proc
                self.embeddings = emb
        
        return ProcessingResponse(processing_results, embedding_result)
    
    async def _extract_content(
        self,
        file_content: bytes,
        filename: str,
        content_type: str,
        document_id: str = None,  # ✅ Added
        kb_id: str = None,        # ✅ Added
        tenant_id: str = None     # ✅ Added
    ) -> Dict[str, Any]:
        """
        Extract content from different file types
        
        Args:
            file_content: File content as bytes
            filename: Original filename
            content_type: MIME type
            document_id: Document ID (for image extraction)
            kb_id: Knowledge base ID (for image extraction)
            tenant_id: Tenant ID (for image extraction)
        """
        file_ext = filename.lower().split('.')[-1]
        
        if file_ext == 'pdf' or 'pdf' in content_type:
            # ✅ Pass IDs to PDF extraction for image handling
            return await self._extract_pdf(
                file_content,
                document_id=document_id,
                kb_id=kb_id,
                tenant_id=tenant_id
            )
        elif file_ext in ['docx', 'doc'] or 'word' in content_type:
            return await self._extract_docx(file_content)
        elif file_ext in ['pptx', 'ppt'] or 'presentation' in content_type:
            return await self._extract_pptx(file_content)
        elif file_ext in ['xlsx', 'xls'] or 'spreadsheet' in content_type:
            return await self._extract_xlsx(file_content)
        elif file_ext == 'txt':
            return await self._extract_txt(file_content)
        elif file_ext == 'html':
            return await self._extract_html(file_content)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    
    
    async def _extract_pdf(
        self, 
        file_content: bytes, 
        document_id: str = None, 
        kb_id: str = None, 
        tenant_id: str = None
    ) -> Dict[str, Any]:
        """
        Extract text AND images from PDF
        
        Args:
            file_content: PDF file bytes
            document_id: Document ID (required for image extraction)
            kb_id: Knowledge base ID (required for image extraction)
            tenant_id: Tenant ID (required for image extraction)
        """
        try:
            import fitz  # PyMuPDF
            from pypdf import PdfReader
            
            # Extract text using PyPDF2
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)
            
            text_parts = []
            total_pages = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text:
                    text_parts.append(f"[Page {page_num}]\n{text}")
            
            full_text = "\n\n".join(text_parts)
            
            # Extract images if document_id, kb_id, and tenant_id are provided
            extracted_images = []
            image_count = 0
            
            if document_id and kb_id and tenant_id:
                try:
                    logger.info(f"Extracting images from PDF document {document_id}...")
                    
                    # Extract images using PyMuPDF
                    extracted_images = await self.pdf_image_extractor.extract_pdf_images(
                        pdf_content=file_content,
                        document_id=document_id,
                        kb_id=kb_id,
                        tenant_id=tenant_id
                    )
                    
                    image_count = len(extracted_images)
                    logger.info(f"✅ Extracted {image_count} images from PDF")
                    
                    # Process each extracted image for embeddings
                    if image_count > 0:
                        await self._process_extracted_images(extracted_images, kb_id, tenant_id)
                        
                except Exception as e:
                    logger.error(f"Error extracting images from PDF: {e}")
                    # Don't fail the whole document processing if image extraction fails
                    image_count = 0
            else:
                logger.warning("Skipping image extraction: document_id, kb_id, or tenant_id not provided")
            
            return {
                "text": full_text,
                "pages": total_pages,
                "images": image_count,
                "extracted_images": extracted_images,  # Image metadata list
                "tables": 0  # TODO: Detect tables
            }
            
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise


    async def _process_extracted_images(
        self, 
        extracted_images: List[Dict[str, Any]], 
        kb_id: str, 
        tenant_id: str
    ):
        """
        Process extracted PDF images: generate embeddings and save to database
        
        Args:
            extracted_images: List of extracted image metadata
            kb_id: Knowledge base ID
            tenant_id: Tenant ID
        """
        try:
            from app.main import get_image_processor
            from app.services.image_vector_store import ImageVectorStore
            import mysql.connector
            from app.config import settings
            from PIL import Image
            
            processor = get_image_processor()
            vector_store = ImageVectorStore(kb_id)
            
            # Database connection
            conn = mysql.connector.connect(
                host=settings.DB_HOST,
                user=settings.DB_USER,
                password=settings.DB_PASSWORD,
                database=settings.DB_NAME
            )
            cursor = conn.cursor()
            
            logger.info(f"Processing {len(extracted_images)} extracted images...")
            
            for img_meta in extracted_images:
                try:
                    # Read image file from disk
                    image_path = img_meta["storage_path"]
                    
                    with open(image_path, "rb") as f:
                        image_bytes = f.read()
                    
                    # Verify file was saved correctly
                    if len(image_bytes) == 0:
                        logger.error(f"Image file is empty: {image_path}")
                        continue
                    
                    # Open image with PIL
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    
                    # Convert to RGB if necessary for CLIP
                    if pil_image.mode != "RGB":
                        pil_image = pil_image.convert("RGB")
                    
                    # Generate embedding
                    embedding_result = await processor.generate_image_embedding(pil_image)
                    
                    # Save to vector store
                    await vector_store.store_image(
                        image_id=img_meta["id"],
                        embedding=embedding_result["embedding"],
                        metadata=img_meta
                    )
                    
                    # Save to MySQL database
                    cursor.execute(
                        """INSERT INTO yovo_tbl_aiva_images (
                            id, kb_id, tenant_id, filename, storage_url, content_type,
                            width, height, file_size_bytes, metadata, created_at
                        ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, NOW())""",
                        (
                            img_meta["id"],
                            kb_id,
                            tenant_id,
                            img_meta["filename"],
                            img_meta["storage_url"],
                            img_meta["content_type"],
                            img_meta["width"],
                            img_meta["height"],
                            img_meta["file_size_bytes"],
                            json.dumps({
                                "document_id": img_meta["document_id"],
                                "page_number": img_meta["page_number"],
                                "image_index": img_meta["image_index"],
                                "embedding_dimension": img_meta["embedding_dimension"]
                            })
                        )
                    )
                    
                    logger.info(f"✅ Processed image {img_meta['id']} from page {img_meta['page_number']}")
                    
                except Exception as e:
                    logger.error(f"Error processing image {img_meta.get('id', 'unknown')}: {e}")
                    continue
            
            conn.commit()
            cursor.close()
            conn.close()
            
            logger.info(f"✅ Successfully processed all extracted images")
            
        except Exception as e:
            logger.error(f"Error in _process_extracted_images: {e}")
            raise
            
    
    async def _store_images_in_db(
        self, 
        image_metadata_list: List[Dict],
        kb_id: str,
        tenant_id: str
    ):
        """Store extracted images in database with embeddings"""
        import mysql.connector
        
        conn = mysql.connector.connect(
            host=settings.DB_HOST,
            port=settings.DB_PORT,
            user=settings.DB_USER,
            password=settings.DB_PASSWORD,
            database=settings.DB_NAME
        )
        cursor = conn.cursor()
        
        try:
            for img_meta in image_metadata_list:
                # Store in yovo_tbl_aiva_images
                cursor.execute("""
                    INSERT INTO yovo_tbl_aiva_images (
                        id, kb_id, tenant_id, filename, storage_url,
                        image_type, width, height, file_size_bytes,
                        description, metadata, vector_id
                    ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                """, (
                    img_meta["image_id"],
                    kb_id,
                    tenant_id,
                    f"page_{img_meta['page_number']}_img_{img_meta['image_index']}.{img_meta['format']}",
                    img_meta["storage_path"],
                    f"image/{img_meta['format']}",
                    img_meta["width"],
                    img_meta["height"],
                    0,  # file_size_bytes - calculate if needed
                    f"Image from page {img_meta['page_number']} of document",
                    json.dumps({
                        "document_id": img_meta["document_id"],
                        "page_number": img_meta["page_number"],
                        "image_index": img_meta["image_index"],
                        "embedding_dimension": img_meta["embedding_dimension"]
                    }),
                    f"clip_{img_meta['image_id']}"
                ))
                
                # Store CLIP embedding in Redis (via ImageVectorStore)
                from app.services.image_vector_store import ImageVectorStore
                image_vector_store = ImageVectorStore(kb_id)
                await image_vector_store.add_image(
                    image_id=img_meta["image_id"],
                    embedding=img_meta["embedding"],
                    metadata={
                        "document_id": img_meta["document_id"],
                        "page_number": img_meta["page_number"],
                        "storage_path": img_meta["storage_path"],
                        "width": img_meta["width"],
                        "height": img_meta["height"]
                    }
                )
            
            conn.commit()
            logger.info(f"Stored {len(image_metadata_list)} images in database")
            
        except Exception as e:
            conn.rollback()
            print(f"Error storing images: {e}")
            raise
        finally:
            cursor.close()
            conn.close()
    
    async def _extract_docx(self, file_content: bytes) -> Dict[str, Any]:
        """Extract text from DOCX with markdown formatting"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)
            
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Preserve heading styles
                    if para.style.name.startswith('Heading'):
                        level = para.style.name[-1]
                        text_parts.append(f"{'#' * int(level)} {para.text}")
                    else:
                        text_parts.append(para.text)
            
            # Extract tables with markdown
            for table in doc.tables:
                table_md = []
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    table_md.append(f"| {row_text} |")
                
                if table_md:
                    text_parts.append("\n" + "\n".join(table_md) + "\n")
            
            full_text = "\n\n".join(text_parts)
            
            return {
                "text": full_text,
                "pages": len(doc.sections),
                "images": 0,
                "tables": len(doc.tables)
            }
            
        except Exception as e:
            print(f"DOCX extraction error: {e}")
            raise
    
    async def _extract_pptx(self, file_content: bytes) -> Dict[str, Any]:
        """Extract text from PPTX with markdown formatting"""
        try:
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"## Slide {slide_num}\n"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))
            
            full_text = "\n\n---\n\n".join(text_parts)
            
            return {
                "text": full_text,
                "pages": len(prs.slides),
                "images": 0,
                "tables": 0
            }
            
        except Exception as e:
            print(f"PPTX extraction error: {e}")
            raise
    
    async def _extract_xlsx(self, file_content: bytes) -> Dict[str, Any]:
        """Extract text from XLSX with markdown table formatting"""
        try:
            xlsx_file = io.BytesIO(file_content)
            wb = load_workbook(xlsx_file, data_only=True)
            
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"## Sheet: {sheet_name}\n"]
                
                # Format as markdown table
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    sheet_text.append(f"| {row_text} |")
                
                if len(sheet_text) > 1:
                    text_parts.append("\n".join(sheet_text))
            
            full_text = "\n\n---\n\n".join(text_parts)
            
            return {
                "text": full_text,
                "pages": len(wb.sheetnames),
                "images": 0,
                "tables": len(wb.sheetnames)
            }
            
        except Exception as e:
            logger.error(f"XLSX extraction error: {e}")
            raise
    
    async def _extract_txt(self, file_content: bytes) -> Dict[str, Any]:
        """Extract text from TXT"""
        try:
            text = file_content.decode('utf-8', errors='ignore')
            
            return {
                "text": text,
                "pages": 1,
                "images": 0,
                "tables": 0
            }
            
        except Exception as e:
            logger.error(f"TXT extraction error: {e}")
            raise
    
    async def _extract_html(self, file_content: bytes) -> Dict[str, Any]:
        """Extract text from HTML with markdown conversion"""
        try:
            html = file_content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html, 'lxml')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Convert HTML structure to markdown-ish format
            text_parts = []
            
            for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                level = int(heading.name[1])
                text_parts.append(f"{'#' * level} {heading.get_text().strip()}")
            
            for para in soup.find_all('p'):
                text = para.get_text().strip()
                if text:
                    text_parts.append(text)
            
            full_text = "\n\n".join(text_parts)
            
            return {
                "text": full_text,
                "pages": 1,
                "images": 0,
                "tables": 0
            }
            
        except Exception as e:
            logger.error(f"HTML extraction error: {e}")
            raise
    
    async def process_text_content(
        self,
        document_id: str,
        kb_id: str,
        tenant_id: str,
        text: str,
        title: str = None,
        source_url: str = None,
        metadata: Dict[str, Any] = None
    ) -> Any:
        """Process plain text content (for web scraping, etc)"""
        import time
        start_time = time.time()
        
        if metadata is None:
            metadata = {}
        
        # Process text with markdown preservation
        processed = await self.text_processor.process_text(
            text=text,
            document_id=document_id,
            kb_id=kb_id,
            metadata={
                **metadata,
                "title": title,
                "source_url": source_url,
                "content_type": "text/html"
            },
            preserve_formatting=True  # NEW
        )
        
        # Generate embeddings
        embeddings_result = await self.embedding_service.generate_embeddings_for_chunks(
            chunks=processed["chunks"]
        )
        
        # Store in vector store
        await self.vector_store.store_document(
            document_id=document_id,
            kb_id=kb_id,
            tenant_id=tenant_id,
            chunks=processed["chunks"],
            embeddings=embeddings_result["embeddings"]
        )
        
        processing_time = int((time.time() - start_time) * 1000)
        
        # Build response
        from app.models.responses import DocumentProcessingResult, EmbeddingResult
        
        processing_results = DocumentProcessingResult(
            total_pages=1,
            total_chunks=len(processed["chunks"]),
            extracted_images=0,
            detected_tables=0,
            detected_faqs=processed.get("faqs", 0),
            chunks_by_type=processed.get("chunks_by_type", {}),
            language_detected=processed.get("languages", []),
            has_roman_urdu=processed.get("has_roman_urdu", False),
            processing_time_ms=processing_time
        )
        
        embedding_result = EmbeddingResult(
            total_embeddings_generated=embeddings_result["total_embeddings"],
            embedding_model=embeddings_result["model"],
            total_tokens_embedded=embeddings_result["total_tokens"],
            vector_dimension=embeddings_result["dimension"]
        )
        
        class ProcessingResponse:
            def __init__(self, proc, emb):
                self.processing_results = proc
                self.embeddings = emb
        
        return ProcessingResponse(processing_results, embedding_result)
    
    
    async def get_document_status(self, document_id: str) -> Dict[str, Any]:
        """Get document processing status"""
        # This would check Redis or DB for status
        return {
            "document_id": document_id,
            "status": "completed",
            "message": "Document processed successfully"
        }
    
    
    async def reprocess_document(self, document_id: str):
        """Reprocess existing document"""
        # TODO: Implement reprocessing logic
        raise NotImplementedError("Reprocess not yet implemented")
    
    
    async def get_similar_documents(self, document_id: str, top_k: int):
        """Get similar documents"""
        # TODO: Implement similarity search
        return []
        
        
    async def delete_document(self, document_id: str):
        """Delete document and all associated data"""
        await self.vector_store.delete_document(document_id)
